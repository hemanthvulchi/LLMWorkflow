import os
import sys
import traceback
from PySide6 import QtWidgets
from PySide6.QtWidgets import QLabel, QTextEdit, QPushButton, QFileDialog, QMessageBox, QFormLayout, QApplication
from node_editor.node import Node
from node_editor.configdialog import ConfigDialog
from utils.display import Display
from utils.llmconnection import LLMConnection
from pptx import Presentation
from pptx.util import Inches
import win32com.client
import pandas as pd
import json

class PwrPointAdvancedProcessor_Node(Node):
    def __init__(self):
        super().__init__()
        self.title_text = "Adv PwrPnt Assistant"
        self.type_text = "APIs in beta - WIP"
        self.set_color(title_color=(190, 0, 0))
        self.pin_output = self.add_pin(name="value", is_output=True)
        self.pin_A = self.add_pin(name="input A", is_output=False)
        self.ex_in_pin=self.add_pin(name="Ex In", is_output=False, execution=True)
        self.pin_output = self.add_pin(name="Ex Out", is_output=True, execution=True)        
        self.build()

        self.config = {
            "user_prompt": "content from slide:",
            "system_prompt": "You are a security risk professional. You will be presented with slide data and then reference information,"
                             "with which you will have to reivew the slide content",
            "output":"",
            "additional_input":"",
            "max_tokens": 64,
            "model": "gpt-3.5-turbo",
            "id": "",
            "object": "",
            "temperature": ""
        }

        # Create a single instance of the ExtendedConfigDialog
        
        
    def init_widget(self):
        self.widget = QtWidgets.QWidget()
        self.widget.setFixedWidth(150)
        layout = QtWidgets.QVBoxLayout()
        layout.setContentsMargins(0, 0, 0, 0)
        self.textbox = QTextEdit()
        self.textbox.setFixedHeight(50)
        self.responsetext = QTextEdit()
        self.responsetext.setFixedHeight(50)
        self.responsetext.setStyleSheet("""
        QTextEdit{
        background: rgb(50, 50, 50); /*background color */
        }
        """)
        self.responsetext.setReadOnly(True)

        self.btn = QtWidgets.QPushButton("Refresh")
        self.btn.clicked.connect(self.btn_refresh)

        self.config_btn = QtWidgets.QPushButton("Configure")
        self.config_btn.clicked.connect(self.show_configuration)

        #layout.addWidget(self.textbox)
        layout.addWidget(self.responsetext)
        layout.addWidget(self.btn)
        layout.addWidget(self.config_btn)
        self.widget.setLayout(layout)

        proxy = QtWidgets.QGraphicsProxyWidget()
        proxy.setWidget(self.widget)
        proxy.setParentItem(self)

        super().init_widget()

    def show_configuration(self):
        #self.config['input1'] = 
        if self.pin_A:
            self.config['additional_input'] = str(self.pin_A.connected_pin.get_data())
            print("[Excel Process] added data:",self.config['additional_input'])
        self.dialog = ExtendedConfigDialog(json.dumps(self.config), QtWidgets.QMainWindow())            
        if self.dialog.exec():
            self.config = json.loads(self.dialog.get_configuration())
            self.responsetext.setText(self.config["output"])
            self.textbox.setText(self.dialog.complete_prompt.toPlainText())
            self.pin_output.set_data(self.responsetext.toPlainText())
        print("Pin output data:", self.pin_output.get_data())

    def btn_refresh(self):
        print("Refreshing data")
        if self.pin_A:
            self.pin_A.set_data(self.pin_A.connected_pin.get_data())
            print("pin A set:",self.pin_A.connected_pin.get_data())
            self.responsetext.setText(str(self.pin_A.connected_pin.get_data()))
            self.config['additional_input'] = str(self.pin_A.connected_pin.get_data())
            print(self.config)
            #if self.dialog.update_input(str(self.pin_A.connected_pin.get_data())):
            #print(self.config)


class ExtendedConfigDialog(ConfigDialog):

    
    def __init__(self, config_json, parent=None):
        super(ExtendedConfigDialog, self).__init__(config_json, parent)

        self.ppt_app = ""
        self.file_label = QLabel('Selected file:')
        self.load_button = QPushButton('Load PowerPoint File', self)
        self.load_button.clicked.connect(self.load_file)
        self.comments_prompt = QTextEdit()
        self.comments_prompt.setMinimumHeight(30)
        self.comments_prompt.setMinimumWidth(200)        
        self.top_layout = QFormLayout()
        self.top_layout.addWidget(self.load_button)
        self.top_layout.addWidget(self.file_label)   
        self.top_layout.addWidget(self.comments_prompt)     
        self.main_layout.addLayout(self.top_layout)
        self.setmainlayout()
        self.setLayout(self.main_layout)
        # Initialize prompts attributes
        self.comments_prompt.setText("Answer in yes/no fashion with one line justification for the folling questions\n"
                                     "Is the Data Accurate and Consistent?\n"
                                     "Is the Data Presented Clearly and Concisely?\n"
                                     "Are Visuals and Graphs Correct and Relevant?\n"
                                     "Are there any grammatical or typos in the slide content?\n"
                                     "Slide content starts"
                                     )
        self.user_prompt.setText("content from slide:")
        self.post_prompt.setText("\nEnd of slide content\n\n Please refer to the below report to review the slide content")

    def get_configuration(self):
        config = json.loads(super().get_configuration())
        config["extra_field1"] = self.extra_field1.text()
        config["extra_field2"] = self.extra_field2.text()
        return json.dumps(config)

    def updatecompleteprompt(self):
        tempString = self.user_prompt.toPlainText() + " [Slide Text]" + self.post_prompt.toPlainText() + "\n" + self.additional_input
        self.complete_prompt.setText(tempString)    
    
    def load_file(self):
        file_dialog = QFileDialog()
        file_dialog.setFileMode(QFileDialog.ExistingFile)
        file_dialog.setNameFilter("PowerPoint files (*.pptx)")

        if file_dialog.exec():
            self.selected_file = file_dialog.selectedFiles()[0]
            self.file_label.setText(f'Selected file: {self.selected_file}')
            self.ppt_app = win32com.client.GetObject(self.selected_file)


        
    def test_api_connection(self):
        slide_number = 1
        try:
            connection = LLMConnection()
            rollingText = ""
            for ppt_slide in self.ppt_app.Slides:
                # Extract slide text
                slide_text = ""
                for shape in ppt_slide.Shapes:
                    if shape.HasTextFrame:
                        slide_text += shape.TextFrame.TextRange.Text
                inputText = self.user_prompt.toPlainText() + str(slide_text) + self.post_prompt.toPlainText() +  "\n" + self.additional_input                        
                response = connection.call_prompt(inputText, self.system_prompt.toPlainText(),self.model.text(),self.max_tokens_slider.value())
                outputText = response.choices[0].message.content
                if ppt_slide.HasNotesPage:
                    notes_page = ppt_slide.NotesPage
                    notes_shape = notes_page.Shapes.Placeholders[2]  # Placeholder for notes text
                    existing_notes = notes_shape.TextFrame.TextRange.Text
                    notes_shape.TextFrame.TextRange.Text = f"{existing_notes}\n\Automated Review:\n{outputText}"
                else:
                    ppt_slide.NotesPage.Shapes.Placeholders[2].TextFrame.TextRange.Text = f"Automated Review:\n{outputText}"
                # Add the comment with the combined text
                inputText = self.comments_prompt.toPlainText() + str(slide_text) + self.post_prompt.toPlainText() +  "\n" + self.additional_input                        
                response_commment = connection.call_prompt(inputText, self.system_prompt.toPlainText(),self.model.text(),self.max_tokens_slider.value())
                outputText_comment = response_commment.choices[0].message.content
                comment = ppt_slide.Comments.Add(
                    Left=100,
                    Top=100,
                    Author="sdf dsdf",
                    AuthorInitials="sdf",
                    Text=str(outputText_comment) # Use the combined text
                )
                rollingText = rollingText + "\n" + str(outputText_comment)
            self.ppt_app.Save()
            self.ppt_app.Close()
 
 
            print("in config dialog")

            Display.show_message_box("SuccessExtend", "API connection successful Extend!\nResponse: " + rollingText)
            self.responseAPItext.setText(rollingText)
        except Exception as e:
            print("Error", f"API connection failed!\nError: {str(e)}")
            print(traceback.format_exc())
            Display.show_message_box("Error", f"API connection failed!\nError: {str(e)}")

    def add_comment(self, slide, text):
        author = "Processor"
        initials = "P"
        left = Inches(0.5)
        top = Inches(0.5)
        slide.notes_slide.notes_text_frame.text += f"\n\n{text}\n\n({author}, {initials})"
    
    def add_note(self, slide, text):
        if slide.notes_slide:
            slide.notes_slide.notes_text_frame.text += f"\n\n{text}"
        else:
            slide.notes_slide = slide.placeholders[1]
            slide.notes_slide.text = text

if __name__ == "__main__":
    app = QtWidgets.QApplication(sys.argv)
    node = PwrPointAdvancedProcessor_Node()
    node.show()
    sys.exit(app.exec())
