from utils import directory as dir
import hashlib
import pickle
import chromadb
import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

#Creating a parent class in case we would want to use a different vector database
class VectorDB:
    def __init__(self):
        # Initialize a connection to the ChromaVectorDB
        self.connection = ChromaVectorDB()

    def query_db(self, *args, **kwargs):
        # Query the database using the provided arguments
        return self.connection.query_chromaDB(*args, **kwargs)

    def query_db_StringContext(self, *args, **kwargs):
        return self.connection.query_chromaDB_getContext(*args,*kwargs)

    def update_db(self, *args, **kwargs):
        # Query the database using the provided arguments
        return self.connection.update_chromaDB(*args, **kwargs)
    
    def reload_db(self, *args, **kwargs):
        # Query the database using the provided arguments
        return self.connection.reload_chromaDB(*args, **kwargs)    

#ChromaDB vector database
class ChromaVectorDB:
    _instance = None

    def __new__(cls, *args, **kwargs):
        if not cls._instance:
            cls._instance = super(ChromaVectorDB, cls).__new__(cls, *args, **kwargs)
            cls._instance._initialized = False
        return cls._instance

    def __init__(self):
        if self._initialized:
            return
        self._initialized = True
        self._connect()
        self.update_chromaDB()

    def _connect(self):
        print("Initiating database connection")
        self.documents_folder = dir.check_directory("documents")
        self.settings_folder = dir.check_directory("settings")
        self.db_folder = dir.check_directory("db")
        db_file_path = dir.get_db_filepath(self.db_folder)

        if os.path.exists(db_file_path):
            self._client = chromadb.PersistentClient(path=self.db_folder)
            self._collection = self._client.get_collection("my_collection")
            print("Existing database loaded successfully.")
        else:
            print("No database found. New database being loaded...")
            self._client = chromadb.PersistentClient(path=self.db_folder)
            self._collection = self._client.create_collection(name="my_collection")

    def calculate_file_hash(self, file_path):
        hasher = hashlib.sha256()
        with open(file_path, 'rb') as f:
            buf = f.read()
            hasher.update(buf)
        return hasher.hexdigest()

    def load_file_list(self, current_directory, filename='file_list.pkl'):
        filepath = os.path.join(current_directory, filename)
        if os.path.exists(filepath):
            with open(filepath, 'rb') as f:
                return pickle.load(f)
        return {}

    def save_file_list(self, file_list, current_directory, filename='file_list.pkl'):
        filepath = os.path.join(current_directory, filename)
        with open(filepath, 'wb') as f:
            pickle.dump(file_list, f)

    def scan_documents_folder(self, folder):
        current_files = {}
        for root, _, files in os.walk(folder):
            for file in files:
                file_path = os.path.join(root, file)
                file_hash = self.calculate_file_hash(file_path)
                current_files[file_path] = file_hash
        return current_files

    def add_document_collection(self, file_path):
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            if file_extension == '.txt':
                with open(file_path, "r", encoding="utf-8") as file:
                    document_text = file.read()
            elif file_extension == '.pdf':
                document_text = self.extract_text_from_pdf(file_path)
            elif file_extension == '.pptx':
                document_text = self.extract_text_from_pptx(file_path)
            elif file_extension == '.docx':
                document_text = self.extract_text_from_docx(file_path)
            else:
                print(f"Unsupported file type: {file_extension}")
                return
            
            self._collection.add(
                documents=[document_text],
                ids=[file_path],
                metadatas=[{"source": file_path}]
            )
            print(f"Added document: {file_path}")
        except FileNotFoundError:
            print(f"File not found: {file_path}")
        except Exception as e:
            print(f"Error adding document {file_path}, error: {e}")

    def extract_text_from_pdf(self, file_path):
        text = ""
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text()
        return text

    def extract_text_from_pptx(self, file_path):
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def extract_text_from_docx(self, file_path):
        text = ""
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def remove_document_collection(self, file_path):
        self._collection.delete(ids=[file_path])
        print(f"Removed document: {file_path}")

    def query_chromaDB(self, query_text, noofresults=10):
        results = self._collection.query(
            query_texts=[query_text],
            n_results=noofresults
        )
        return results

    def query_chromaDB_getContext(self,query_text, noofresults=10):
        results = self.query_chromaDB(query_text, noofresults)
        context = "\n".join(results["documents"][0])
        #retrieved_documents = [result["document"] for result in results['documents'][0]]
        #context = f"Query: {query_text}\nRelevant Documents:\n" + "\n".join(retrieved_documents)
        return context


    def update_chromaDB(self):
        file_list = self.load_file_list(self.settings_folder)
        current_files = self.scan_documents_folder(self.documents_folder)

        for file_path, file_hash in current_files.items():
            if file_path not in file_list or file_list[file_path] != file_hash:
                file_list[file_path] = file_hash
                self.add_document_collection(file_path)

        for file_path in list(file_list.keys()):
            if file_path not in current_files:
                del file_list[file_path]
                self.remove_document_collection(file_path)

        self.save_file_list(file_list, self.settings_folder)
        print("List of files in vectorDB")
        for file_path in file_list:
            print(file_path, dir.get_current_directory())

    #need to check if the collection should be recreated
    def reload_chromaDB(self):
        self.reset_chromaDB()
        self._client = chromadb.PersistentClient(path=self.db_folder)
        self._collection = self._client.create_collection(name="my_collection")        
        self.update_chromaDB(self)
    
    def reset_chromaDB(self):
        self._client.reset()